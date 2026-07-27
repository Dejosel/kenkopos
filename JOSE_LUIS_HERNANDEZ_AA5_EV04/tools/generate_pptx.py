from pptx import Presentation
from pptx.util import Pt, Inches

slides = [
    {"title": "KENKOPOS — Presentación para KENKO", "bullets": [
        "Resumen ejecutivo del módulo POS y catálogo de productos",
        "Estado actual del desarrollo (CRUD y POS visual)",
        "Siguientes pasos propuestos para producción"
    ], "notes": "Introducción breve: propósito de la presentación y público objetivo (gerencia KENKO)."},

    {"title": "Objetivo del proyecto", "bullets": [
        "Digitalizar gestión de productos para restaurantes rápidos",
        "Proveer una pantalla POS para tomar comanda y cobrar",
        "Servir como base para escalar a sistema transaccional"
    ], "notes": "Explicar valor para operaciones: velocidad de servicio y control de catálogo."},

    {"title": "Estructura del repositorio", "bullets": [
        "Stack PHP puro (v1): `app/`, `public/` — POS y CRUD operativo",
        "Stack Laravel (v2): `laravel-app/` — migración y mejoras",
        "Documentación y scripts SQL en `database/` y `docs/`"
    ], "notes": "Mencionar coexistencia temporal y migración planeada a Laravel."},

    {"title": "Tecnologías usadas", "bullets": [
        "Backend: PHP 8 (PDO) — MySQL con fallback SQLite",
        "Frontend: HTML5, Bootstrap 5, JavaScript (POS interactividad)",
        "Framework: Laravel 13 (migración), Eloquent y Blade"
    ], "notes": "Aclarar diferencias entre versiones v1 y v2 y por qué ambas existen."},

    {"title": "Funcionalidades implementadas", "bullets": [
        "CRUD completo de productos (crear, listar, editar, eliminar)",
        "Interfaz POS interactiva (ticket en memoria, descuentos, impresión)",
        "Inicialización de BD y seeds automáticos (productos demo)"
    ], "notes": "Mostrar evidencias con rutas y archivos clave del repo."},

    {"title": "Funcionalidades no implementadas", "bullets": [
        "Persistencia de ventas/comandas en BD (orders, items)",
        "Gestión de inventarios, clientes y roles de usuario",
        "Reportes y cierre de caja automatizado"
    ], "notes": "Priorizar qué falta para convertirlo en POS productivo."},

    {"title": "CRUD de Productos — Detalle técnico", "bullets": [
        "Modelo: `app/Models/Product.php` (Prepared Statements PDO)",
        "Controlador: `app/Controllers/ProductController.php` (store, update, destroy)",
        "Vistas: `public/products/*.php` y equivalentes Blade en Laravel"
    ], "notes": "Explicar flujo: formulario → controlador → modelo → BD → redirección."},

    {"title": "Tablas MySQL detectadas", "bullets": [
        "`products` (campos: product_id, name, sku, price, category, color, created_at)",
        "En Laravel: tablas auxiliares (`users`, `sessions`, `jobs`, `cache`)",
        "Scripts SQL: `database/kenkopos.sql`, `database/kenkopos_infinityfree.sql`"
    ], "notes": "Indicar que producto es la tabla principal para el MVP POS."},

    {"title": "Historias de Usuario reales", "bullets": [
        "HU-01 Registrar producto — Formulario y validación",
        "HU-02 Visualizar inventario — Listado y búsquedas",
        "HU-03 Editar / HU-04 Eliminar — Flujo y confirmaciones"
    ], "notes": "Remarcar que estas HU están implementadas y documentadas en `docs/`."},

    {"title": "Arquitectura del sistema", "bullets": [
        "MVC simplificado en PHP + patrón Singleton para BD",
        "Migración a Laravel con Eloquent (mejor mantenibilidad)",
        "Frontend POS desacoplado (JS) que consume catálogo en JSON"
    ], "notes": "Sugerir un diagrama de componentes para la siguiente fase."},

    {"title": "UML sugerido", "bullets": [
        "Diagrama de Casos de Uso: Administrador CRUD, Cajero POS",
        "Diagrama de Clases: Database, Product, ProductController, Response",
        "Diagrama de Secuencia: Crear producto y flujo de venta"
    ], "notes": "Ofrecer generar imágenes UML si lo desean (plantuml o draw.io)."},

    {"title": "Mejoras futuras — Prioridad Alta", "bullets": [
        "Unificar código en Laravel y alinear esquema (`category`, `color`)",
        "Persistir ventas: `orders`, `order_items`, `payments`, `tables`",
        "Implementar autenticación y roles (admin, cajero, mesero)"
    ], "notes": "Explicar impacto en operación y tiempos estimados por fase."},

    {"title": "Mejoras futuras — Prioridad Media/Baja", "bullets": [
        "Reportes y dashboards de ventas diarias y producto TOP",
        "Integración con impresoras térmicas y TPV físico",
        "APIs para integraciones con delivery y contabilidad"
    ], "notes": "Plantear roadmap y requisitos no funcionales (seguridad, backup)."},

    {"title": "Roadmap propuesto", "bullets": [
        "Fase 1 (2-4 semanas): Consolidar Laravel + paridad funcional",
        "Fase 2 (4-6 semanas): Persistencia de ventas y roles",
        "Fase 3: Reportes, despliegue y hardening para producción"
    ], "notes": "Incluir estimaciones aproximadas y dependencias técnicas."},

    {"title": "Cierre ejecutivo", "bullets": [
        "KenkoPOS ya tiene CRUD y POS visual operativo",
        "Sigue trabajo para convertirlo en POS transaccional completo",
        "Propuesta: comenzar migración completa a Laravel y persistencia ventas"
    ], "notes": "Invitar a la gerencia a aprobar roadmap y recursos necesarios."}
]

prs = Presentation()
blank_slide_layout = prs.slide_layouts[1]

for s in slides:
    slide = prs.slides.add_slide(blank_slide_layout)
    title = slide.shapes.title
    title.text = s['title']

    # add content textbox
    left = Inches(0.5)
    top = Inches(1.6)
    width = Inches(9)
    height = Inches(4.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(s['bullets']):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = "• " + b
            p.font.size = Pt(18)
        else:
            p = tf.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(18)

    # Notes
    notes_slide = slide.notes_slide
    notes_text_frame = notes_slide.notes_text_frame
    notes_text_frame.text = s['notes']

out_path = 'docs/KENKOPOS_Presentation_KENKO.pptx'
prs.save(out_path)
print('GENERATED:', out_path)
